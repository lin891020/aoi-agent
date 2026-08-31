"""Build the project-journey deck from `scripts/deck_content.py`.

    uv run --with python-pptx --with playwright python scripts/build_deck.py
    uv run --with python-pptx --with playwright python scripts/build_deck.py --embed-video

Three outputs from one source, so they cannot disagree:

* `docs/deck/aoi-agent-journey.zh-TW.pptx` -- 16:9, speaker notes on every
  slide (a plain-language line, the script, and the interviewer's questions).
* `docs/deck/aoi-agent-journey.zh-TW.html` -- the same slides as a page: arrow
  keys to move, `q` for self-test mode (title and questions only), `a` to
  reveal the answers, `n` to fold the notes.
* `docs/deck/study-guide.zh-TW.md` -- the same content as a printable Q&A.

Images are rendered into `docs/deck/img/` on the way: the architecture page
and the two flow diagrams through Playwright, the operating-point curve from
`models/test_predictions.npz` through matplotlib (skipped with a note when
the file is not there), the station screenshots composed from
`docs/screenshots/`.

`--embed-video` writes a second .pptx to `docs/demo/build/` with the demo
recording embedded; that directory is git-ignored, as the recording is.

python-pptx and playwright are not project dependencies -- pass them with
`--with`. Without them the script says so and stops.
"""

from __future__ import annotations

import argparse
import html
import json
import re
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))
sys.path.insert(0, str(ROOT / "src"))

from deck_content import CORE_TOTAL, FIVE_LABELS, LAYERS, SLIDES, Slide, by_layer  # noqa: E402

OUT = ROOT / "docs" / "deck"
IMG = OUT / "img"
SCREENSHOTS = ROOT / "docs" / "screenshots"
DIAGRAMS = ROOT / "docs" / "diagrams"
ARCH_PAGE = ROOT / "docs" / "architecture-diagram.html"
# Hand-authored diagram pages, dark so a slide takes the screenshot unaltered.
AUTHORED = {
    "planner_reach.png": OUT / "planner-reach.zh-TW.html",
    "model_blind_spot.png": OUT / "model-blind-spot.zh-TW.html",
}
PREDICTIONS = ROOT / "models" / "test_predictions.npz"
VIDEO = ROOT / "docs" / "demo" / "aoi-agent-demo-zh.mp4"
VIDEO_LINK = "https://github.com/lin891020/aoi-agent/releases/tag/demo-2026-08-28"
EMBEDDED_OUT = ROOT / "docs" / "demo" / "build" / "aoi-agent-journey.zh-TW.embedded.pptx"

# One palette, both outputs. Dark ground, one accent, a warm mark for mistakes.
INK = "E8E6E1"
DIM = "9AA0A6"
GROUND = "15171A"
PANEL = "1E2125"
ACCENT = "6FC7CF"
WARN = "E2B45F"
FONT = "Noto Sans TC"
MONO = "IBM Plex Mono"


# ----------------------------------------------------------------- images

def render_images(need_playwright: bool = True) -> dict[str, Path]:
    IMG.mkdir(parents=True, exist_ok=True)
    # Start from what is already on disk: the station screenshots come from
    # scripts/deck_screenshots.py, not from here, and a build that returned only
    # its own renders dropped them from three slides without saying so.
    made: dict[str, Path] = {p.name: p for p in IMG.glob("*.png")}
    made.update(_render_pages() if need_playwright else {})
    curve = _render_operating_point()
    if curve:
        made["operating_point.png"] = curve
    screens = _compose_screens()
    if screens:
        made["screens.png"] = screens
    for name, fn in (("training_curves.png", _render_training_curves),
                     ("crop_curves.png", _render_crop_curves),
                     ("gate_curves.png", _render_gate_curves)):
        path = fn()
        if path:
            made[name] = path
    return made


def _render_pages() -> dict[str, Path]:
    try:
        from playwright.sync_api import sync_playwright
    except ImportError:
        print("playwright not installed: architecture.png and flows.png not rendered "
              "(pass --with playwright)")
        return {}
    made: dict[str, Path] = {}
    with sync_playwright() as p:
        browser = p.chromium.launch()
        # The architecture page: only its diagram, not the page's margins --
        # on a slide the margins were most of the picture.
        page = browser.new_page(viewport={"width": 1600, "height": 1000}, device_scale_factor=2)
        page.goto(ARCH_PAGE.as_uri())
        page.wait_for_load_state("networkidle")
        page.wait_for_timeout(800)
        target = IMG / "architecture.png"
        svg = page.locator("svg").first
        if svg.count():
            svg.screenshot(path=str(target))
        else:
            page.screenshot(path=str(target), full_page=True)
        made["architecture.png"] = target
        # One flow per picture, each at full width of its own slide.
        for name, file in (("flow_disposition.png", "disposition-flow-light.zh-TW.svg"),
                           ("flow_analysis.png", "analysis-flow-light.zh-TW.svg")):
            html_path = IMG / f"_{name}.html"
            html_path.write_text("<!doctype html><meta charset='utf-8'><body style='margin:0;"
                                 "background:#f5f5f5;padding:16px;width:1400px;box-sizing:border-box'>"
                                 + (DIAGRAMS / file).read_text() + "</body>")
            page = browser.new_page(viewport={"width": 1400, "height": 900}, device_scale_factor=2)
            page.goto(html_path.as_uri())
            page.wait_for_timeout(600)
            target = IMG / name
            page.locator("svg").first.screenshot(path=str(target))
            made[name] = target
            html_path.unlink(missing_ok=True)
        # A short viewport lets full_page cut each authored page at its content.
        for name, source in AUTHORED.items():
            if not source.exists():
                continue
            page = browser.new_page(viewport={"width": 1900, "height": 400},
                                    device_scale_factor=2)
            page.goto(source.as_uri())
            page.wait_for_load_state("networkidle")
            # The page's own headline is for the standalone read; on a slide the
            # slide already has a title and two headlines is one too many.
            page.add_style_tag(content="header{display:none}.wrap{padding-top:26px}")
            page.wait_for_timeout(700)
            target = IMG / name
            page.screenshot(path=str(target), full_page=True)
            made[name] = target
        browser.close()
    for name in ("architecture.png", "flow_disposition.png", "flow_analysis.png"):
        if name in made:
            _trim(made[name])
    return made


def _trim(path: Path, margin: int = 24) -> None:
    """Cut a rendered diagram down to its ink plus a small margin.

    The pages carry their own padding inside the element, and on a slide that
    padding was a third of the picture -- the diagram everyone came to see
    sat small in the middle of a grey box.
    """
    from PIL import Image, ImageChops
    im = Image.open(path).convert("RGB")
    bg = Image.new("RGB", im.size, im.getpixel((2, 2)))
    diff = ImageChops.difference(im, bg).convert("L").point(lambda v: 255 if v > 12 else 0)
    box = diff.getbbox()
    if not box:
        return
    left, top, right, bottom = box
    im.crop((max(0, left - margin), max(0, top - margin),
             min(im.width, right + margin), min(im.height, bottom + margin))).save(path)


def _cjk_font(matplotlib) -> None:
    """A face that has the CJK glyphs, whichever of these the machine holds."""
    from matplotlib import font_manager
    have = {f.name for f in font_manager.fontManager.ttflist}
    for name in ("Noto Sans TC", "PingFang TC", "Heiti TC", "Hiragino Sans", "Arial Unicode MS"):
        if name in have:
            matplotlib.rcParams["font.family"] = [name, "DejaVu Sans"]
            return


def _render_operating_point() -> Path | None:
    if not PREDICTIONS.exists():
        print(f"{PREDICTIONS.relative_to(ROOT)} not found: operating_point.png not rendered "
              "(run scripts/train.py)")
        return None
    import numpy as np
    import matplotlib
    matplotlib.use("Agg")
    _cjk_font(matplotlib)
    import matplotlib.pyplot as plt
    from aoi_agent.vision.operating_point import best_at_escape_budget, sweep

    data = np.load(PREDICTIONS)
    names = [str(n) for n in data["label_names"]]
    fc = names.index("false_call")
    points = sweep(data["probabilities"][:, fc], data["labels"], fc)
    best = best_at_escape_budget(points, 0.005)
    xs = [p.escape_rate * 100 for p in points]
    ys = [p.review_reduction * 100 for p in points]

    fig, ax = plt.subplots(figsize=(9, 5), dpi=160, facecolor="#" + PANEL)
    ax.set_facecolor("#" + PANEL)
    ax.plot(xs, ys, color="#" + ACCENT, linewidth=2)
    ax.axvline(0.5, color="#" + WARN, linestyle="--", linewidth=1.2)
    if best:
        ax.scatter([best.escape_rate * 100], [best.review_reduction * 100],
                   color="#" + WARN, zorder=5, s=50)
        ax.annotate(f"門檻 {best.threshold:.3f}\n{best.review_reduction * 100:.1f}% 省掉 @ "
                    f"{best.escape_rate * 100:.2f}% 漏網",
                    (best.escape_rate * 100, best.review_reduction * 100),
                    xytext=(1.2, best.review_reduction * 100 - 18), color="#" + INK,
                    fontsize=11, arrowprops={"color": "#" + DIM, "arrowstyle": "->"})
    ax.set_xlim(0, 3)
    ax.set_ylim(0, 100)
    ax.set_xlabel("escape rate（漏網率，% of defects）", color="#" + INK)
    ax.set_ylabel("review removed（省掉的人工，% of candidates）", color="#" + INK)
    ax.tick_params(colors="#" + DIM)
    for spine in ax.spines.values():
        spine.set_color("#" + DIM)
    ax.grid(color="#33383D", linestyle=":", linewidth=0.8)
    ax.text(0.55, 3, "QP-110 預算 0.5%", color="#" + WARN, fontsize=10)
    fig.tight_layout()
    target = IMG / "operating_point.png"
    fig.savefig(target, facecolor=fig.get_facecolor())
    plt.close(fig)
    return target


def _dark_axes(ax):
    ax.set_facecolor("#" + PANEL)
    ax.tick_params(colors="#" + DIM)
    for spine in ax.spines.values():
        spine.set_color("#" + DIM)
    ax.grid(color="#33383D", linestyle=":", linewidth=0.8)
    ax.xaxis.label.set_color("#" + INK)
    ax.yaxis.label.set_color("#" + INK)
    ax.title.set_color("#" + INK)


def _history_plot(history_path: Path, target: Path, title: str) -> Path | None:
    if not history_path.exists():
        print(f"{history_path.relative_to(ROOT)} not found: {target.name} not rendered")
        return None
    import json
    import matplotlib
    matplotlib.use("Agg")
    _cjk_font(matplotlib)
    import matplotlib.pyplot as plt
    rows = json.loads(history_path.read_text())
    epochs = [r["epoch"] for r in rows]
    fig, axes = plt.subplots(1, 3, figsize=(12, 3.8), dpi=160, facecolor="#" + PANEL)
    for ax, key, label, colour in zip(
        axes, ("train_loss", "val_accuracy", "val_review_reduction"),
        ("train loss", "val accuracy", "val review reduction @ ≤0.5% 漏網"),
        (DIM, DIM, ACCENT),
    ):
        ys = [r[key] for r in rows]
        ax.plot(epochs, ys, marker="o", color="#" + colour, linewidth=2)
        ax.set_title(label, fontsize=11)
        ax.set_xlabel("epoch")
        _dark_axes(ax)
        if key != "train_loss":
            ax.set_ylim(0, 1)
        best = max(range(len(ys)), key=lambda i: ys[i]) if key == "val_review_reduction" else None
        if best is not None:
            ax.scatter([epochs[best]], [ys[best]], color="#" + WARN, zorder=5, s=45)
            ax.annotate(f"{ys[best]:.3f} @ epoch {epochs[best]}", (epochs[best], ys[best]),
                        xytext=(0, 10), textcoords="offset points", ha="center",
                        color="#" + WARN, fontsize=9)
    fig.suptitle(title, color="#" + INK, fontsize=12)
    fig.tight_layout()
    fig.savefig(target, facecolor=fig.get_facecolor())
    plt.close(fig)
    return target


def _render_training_curves() -> Path | None:
    return _history_plot(ROOT / "models" / "history.json", IMG / "training_curves.png",
                         "主線複判模型 · ResNet-18 三通道 · 10 epochs")


def _render_crop_curves() -> Path | None:
    return _history_plot(ROOT / "models" / "pcbaoi_reverifier" / "history.json",
                         IMG / "crop_curves.png", "無樣板 crop 複判器 · 同一個 ResNet-18 · 10 epochs")


def _render_gate_curves() -> Path | None:
    import json
    files = {
        "DeepPCB（二值圖，trainval）": ROOT / "eval" / "results" / "gate_check.json",
        "HRIPCB 對齊（照片）": ROOT / "eval" / "results" / "gate_check_hripcb_aligned.json",
        "HRIPCB 對齊＋擾動": ROOT / "eval" / "results" / "gate_check_hripcb_aligned_perturbed.json",
    }
    if not all(f.exists() for f in files.values()):
        print("eval/results/gate_check*.json missing: gate_curves.png not rendered")
        return None
    import matplotlib
    matplotlib.use("Agg")
    _cjk_font(matplotlib)
    import matplotlib.pyplot as plt
    fig, (left, right) = plt.subplots(1, 2, figsize=(12, 4.2), dpi=160, facecolor="#" + PANEL)
    colours = [ACCENT, WARN, "FF8F7C"]
    for (label, path), colour in zip(files.items(), colours):
        runs = json.loads(path.read_text())["runs"]
        thr = [r["threshold"] for r in runs]
        left.plot(thr, [r["recall"] * 100 for r in runs], marker="o", color="#" + colour, label=label)
        right.plot(thr, [r["mean_false_calls"] for r in runs], marker="o", color="#" + colour, label=label)
    left.set_xlabel("灰階門檻"); left.set_ylabel("recall（%）"); left.set_ylim(0, 100)
    right.set_xlabel("灰階門檻"); right.set_ylabel("每張圖的平均誤報數"); right.set_yscale("log")
    for ax in (left, right):
        _dark_axes(ax)
        ax.axvline(60, color="#" + DIM, linestyle="--", linewidth=1)
        ax.legend(facecolor="#" + PANEL, edgecolor="#33383D", labelcolor="#" + INK, fontsize=9)
    left.set_title("門檻 60 是主線的設定", fontsize=11)
    right.set_title("對數尺度", fontsize=11)
    fig.tight_layout()
    target = IMG / "gate_curves.png"
    fig.savefig(target, facecolor=fig.get_facecolor())
    plt.close(fig)
    return target


def _compose_screens() -> Path | None:
    files = [SCREENSHOTS / f"{n}-zh.png" for n in ("queue", "region", "boards", "ask")]
    if not all(f.exists() for f in files):
        print("docs/screenshots/*-zh.png missing: screens.png not rendered")
        return None
    import matplotlib
    matplotlib.use("Agg")
    _cjk_font(matplotlib)
    import matplotlib.image as mpimg
    import matplotlib.pyplot as plt

    fig, axes = plt.subplots(2, 2, figsize=(12, 7.2), dpi=140, facecolor="#" + GROUND)
    for ax, f, label in zip(axes.flat, files, ("待複判", "區域頁", "PCB 處置", "產線查詢")):
        img = mpimg.imread(f)
        h = min(img.shape[0], int(img.shape[1] * 0.62))
        ax.imshow(img[:h])
        ax.set_title(label, color="#" + INK, fontsize=11, loc="left")
        ax.axis("off")
    fig.tight_layout()
    target = IMG / "screens.png"
    fig.savefig(target, facecolor=fig.get_facecolor())
    plt.close(fig)
    return target


# ------------------------------------------------------------------ pptx

def _typeface(run, latin: str, cjk: str = FONT) -> None:
    """Name the East Asian typeface as well as the Latin one.

    `run.font.name` writes only `<a:latin>`. Chinese glyphs then fall to the
    theme's East Asian default: PowerPoint quietly substitutes something that
    has them, LibreOffice substitutes something that does not and every CJK
    character in the run vanishes. Setting `<a:ea>` explicitly is what makes
    the file render the same on a machine that is not this one -- and it is
    always the sans face, so a run set in the mono face keeps mono digits and
    readable Chinese.
    """
    from pptx.oxml.ns import qn

    run.font.name = latin
    rPr = run._r.get_or_add_rPr()
    anchor = rPr.find(qn("a:latin"))
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            anchor.addnext(el)
        el.set("typeface", cjk)
        anchor = el



def _notes_text(slide: Slide) -> str:
    lines = [f"白話：{slide.plain}", ""]
    if slide.hero and slide.bullets:
        lines += ["頁面上只有圖；要點："] + [f"• {b}" for b in slide.bullets] + [""]
    lines += [f"講稿 {i}. {s}" for i, s in enumerate(slide.notes, 1)]
    if slide.questions:
        lines += ["", "面試官會問："]
        for q, a in slide.questions:
            lines += [f"Q：{q}", f"A：{a}", ""]
    return "\n".join(lines).rstrip()


def build_pptx(target: Path, images: dict[str, Path], embed_video: bool) -> Path:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Emu, Inches, Pt

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    W, H = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]
    rgb = lambda hexa: RGBColor.from_string(hexa)  # noqa: E731

    def text(slide, left, top, width, height, runs, size=16, color=INK, bold=False,
             align=None, font=FONT, spacing=1.15):
        """`runs`: list of paragraphs; a paragraph is a str or a list of
        (text, {bold, color, size, font}) tuples."""
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        first = True
        for para in runs:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.line_spacing = spacing
            if align:
                p.alignment = align
            parts = [(para, {})] if isinstance(para, str) else para
            for chunk, style in parts:
                r = p.add_run()
                r.text = chunk
                _typeface(r, style.get("font", font))
                r.font.size = Pt(style.get("size", size))
                r.font.bold = style.get("bold", bold)
                r.font.color.rgb = rgb(style.get("color", color))
        return box

    FIGURE = re.compile(r"\*\*(.+?)\*\*|(\d[\d,./%×x−\-–→ ]*\d%?|\d%?)")

    def runs(body: str, size: float, colour: str = INK):
        """Paragraph runs for one cell: `**...**` and every figure in bold accent.

        A wall of 14pt prose has no landmarks; the numbers are what the
        reader must take away, so they are the landmarks.
        """
        out, at = [], 0
        for m in FIGURE.finditer(body):
            if m.start() > at:
                out.append((body[at:m.start()], {"size": size, "color": colour}))
            text_of = m.group(1) or m.group(2)
            out.append((text_of, {"size": size, "color": ACCENT, "bold": True}))
            at = m.end()
        if at < len(body):
            out.append((body[at:], {"size": size, "color": colour}))
        return out or [(body, {"size": size, "color": colour})]

    def ground(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = rgb(GROUND)

    def chrome(slide, s: Slide, number: int):
        text(slide, Inches(0.5), Inches(7.0), Inches(6), Inches(0.35),
             [[(s.layer, {"color": DIM, "size": 10, "font": MONO})]])
        tail = f"{number} / {len(SLIDES)}"
        if s.core:
            tail = f"★ 主線 {s.core}/{CORE_TOTAL} · " + tail
        text(slide, Inches(8.8), Inches(7.0), Inches(4.0), Inches(0.35),
             [[(tail, {"color": ACCENT if s.core else DIM, "size": 10, "font": MONO})]],
             align=PP_ALIGN.RIGHT)

    def title_line(slide, s: Slide, size=28):
        """The title, and its clause after the first comma as a dim subtitle.

        「錯 #1：整線漏網率 5.4% → 0.61%，我算錯了，而且錯在對我有利的方向」is
        a sentence, not a title; the hook keeps its place, one size down.
        Returns the y where the body may start.
        """
        head, tail = s.title, ""
        if len(s.title) > 22 and "，" in s.title:
            head, tail = s.title.split("，", 1)
        text(slide, Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8),
             [[(head, {"bold": True, "size": size})]])
        if tail:
            text(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(0.45),
                 [[(tail, {"size": 15, "color": DIM})]])
            return Inches(1.45)
        return Inches(1.25)

    def bullets(slide, items, left, top, width, height, size=15):
        text(slide, left, top, width, height,
             [[("• ", {"color": ACCENT, "size": size})] + runs(b, size) for b in items],
             spacing=1.25)

    def stat_strip(slide, stats, top):
        """Two or three numbers to remember, each in its own tile, one row."""
        n = len(stats)
        tile_w = (W - Inches(1.0) - Inches(0.25) * (n - 1)) / n
        tile_h = Inches(0.95)
        for i, (head, label) in enumerate(stats):
            x = Inches(0.5) + i * (tile_w + Inches(0.25))
            box = slide.shapes.add_shape(1, x, top, tile_w, tile_h)
            box.fill.solid(); box.fill.fore_color.rgb = rgb(PANEL)
            box.line.color.rgb = rgb("33383D")
            box.shadow.inherit = False
            size = 22 if len(head) <= 9 else (18 if len(head) <= 14 else 15)
            text(slide, x + Inches(0.2), top + Inches(0.08), tile_w - Inches(0.4), Inches(0.5),
                 [[(head, {"bold": True, "size": size, "color": ACCENT, "font": MONO})]])
            text(slide, x + Inches(0.2), top + Inches(0.55), tile_w - Inches(0.4), Inches(0.35),
                 [[(label, {"size": 10.5, "color": DIM})]])
        return top + tile_h + Inches(0.18)

    def five_box(slide, s: Slide, left, top, width, height):
        # Rows take the height their text needs; the type is the largest of
        # 15/14/13/12pt at which the five cells fit the height left under the
        # title and the stat strip.
        rows = list(FIVE_LABELS.items())
        gap = Inches(0.14)

        def layout(size):
            cpl = {15: 50, 14: 54, 13: 58, 12: 63}[size]
            line_h = {15: 0.31, 14: 0.29, 13: 0.27, 12: 0.25}[size]
            heights = [Inches(line_h * max(1, -(-len(s.five[k]) // cpl)) + 0.08) for k, _ in rows]
            return heights, sum(heights) + gap * (len(rows) - 1)

        for size in (15, 14, 13, 12):
            heights, total = layout(size)
            if total <= height:
                break
        y = top
        for (key, label), row_h in zip(rows, heights):
            colour = WARN if key == "mistake" else ACCENT
            if key in ("mistake", "rule"):
                # Two colours for the whole deck: the problem sits on a warm
                # band, the answer on the accent band. Nothing else is tinted.
                band = slide.shapes.add_shape(1, left - Inches(0.15), y - Inches(0.04),
                                              width + Inches(0.3), row_h + Inches(0.08))
                band.fill.solid()
                band.fill.fore_color.rgb = rgb("2B2416" if key == "mistake" else "16292B")
                band.line.fill.background()
                band.shadow.inherit = False
            text(slide, left, y, Inches(1.45), row_h,
                 [[(label, {"bold": True, "color": colour, "size": 13})]])
            text(slide, left + Inches(1.5), y, width - Inches(1.5), row_h,
                 [runs(s.five[key], size)], spacing=1.12)
            y += row_h + gap

    for number, s in enumerate(SLIDES, 1):
        slide = prs.slides.add_slide(blank)
        ground(slide)
        if s.kind == "title":
            text(slide, Inches(0.8), Inches(2.3), Inches(11.7), Inches(1.4),
                 [[(s.title, {"bold": True, "size": 40})]])
            text(slide, Inches(0.8), Inches(3.9), Inches(11.7), Inches(1.6),
                 [[(b, {"color": DIM, "size": 16})] for b in s.bullets])
        elif s.kind == "headline":
            title_line(slide, s)
            # Five numbers a manager reads in thirty seconds: the figure large,
            # its sentence small, one tile each. Bullets of prose at 17pt left
            # two thirds of the page empty and the numbers buried mid-sentence.
            tiles = []
            for b in s.bullets:
                head, _, rest = b.partition(" ")
                tiles.append((head, rest))
            n = len(tiles)
            cols = 3 if n > 4 else n
            tile_w = (W - Inches(1.2) - Inches(0.3) * (cols - 1)) / cols
            tile_h = Inches(2.4)
            for i, (head, rest) in enumerate(tiles):
                r, c = divmod(i, cols)
                x = Inches(0.6) + c * (tile_w + Inches(0.3))
                y = Inches(1.5) + r * (tile_h + Inches(0.3))
                box = slide.shapes.add_shape(1, x, y, tile_w, tile_h)  # 1 = rectangle
                box.fill.solid(); box.fill.fore_color.rgb = rgb(PANEL)
                box.line.color.rgb = rgb("33383D")
                box.shadow.inherit = False
                text(slide, x + Inches(0.25), y + Inches(0.2), tile_w - Inches(0.5), Inches(1.0),
                     [[(head, {"bold": True, "size": 34, "color": ACCENT, "font": MONO})]])
                text(slide, x + Inches(0.25), y + Inches(1.15), tile_w - Inches(0.5), tile_h - Inches(1.25),
                     [[(rest, {"size": 13, "color": INK})]], spacing=1.2)
        elif s.kind in ("text", "closing"):
            top = title_line(slide, s) + Inches(0.15)
            if s.table:
                rows, cols = len(s.table), len(s.table[0])
                row_h = Inches(0.5)
                shape = slide.shapes.add_table(rows, cols, Inches(0.6), top, Inches(12.1), row_h * rows)
                tbl = shape.table
                for r, row in enumerate(s.table):
                    for c, cell_text in enumerate(row):
                        cell = tbl.cell(r, c)
                        cell.text = ""
                        para = cell.text_frame.paragraphs[0]
                        run = para.add_run()
                        run.text = cell_text
                        _typeface(run, FONT)
                        run.font.size = Pt(13 if r else 11.5)
                        run.font.bold = r == 0
                        run.font.color.rgb = rgb(ACCENT if r == 0 else INK)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = rgb(PANEL if r % 2 else GROUND)
                        cell.margin_top = cell.margin_bottom = Inches(0.04)
                top = top + row_h * rows + Inches(0.25)
            bullets(slide, s.bullets, Inches(0.6), top, Inches(12.1), H - top - Inches(0.6),
                    size=14 if s.table else 18)
        elif s.kind in ("image", "video"):
            title_line(slide, s)
            img = images.get(s.image or "")

            def place(path, left, top, max_w, max_h):
                pic = slide.shapes.add_picture(str(path), left, top, width=max_w)
                if pic.height > max_h:
                    ratio = max_h / pic.height
                    pic.height = max_h
                    pic.width = Emu(int(pic.width * ratio))
                return pic

            if s.hero and img:
                # The diagram is the slide: as large as the page allows,
                # centred, nothing beside it. The bullets that used to sit at
                # its right are in the notes, where the speaker reads them.
                pic = place(img, Inches(0.4), Inches(1.1), Inches(12.5), Inches(5.3))
                pic.left = Emu(int((W - pic.width) / 2))
                # A picture wider than the frame leaves its slack at the bottom,
                # which reads as an unfinished slide rather than a margin.
                pic.top = Emu(int(Inches(1.1) + (Inches(5.3) - pic.height) / 2))
            elif s.wide and img:
                # The picture across the whole width, the bullets in two
                # columns under it -- a flow diagram at half width was a
                # thumbnail nobody in the room could read.
                pic = place(img, Inches(0.6), Inches(1.3), Inches(12.1), Inches(3.9))
                pic.left = Emu(int((W - pic.width) / 2))
                y = pic.top + pic.height + Inches(0.15)
                half = len(s.bullets) - len(s.bullets) // 2
                bullets(slide, s.bullets[:half], Inches(0.6), y, Inches(5.95), Inches(6.35) - y, size=12)
                bullets(slide, s.bullets[half:], Inches(6.75), y, Inches(5.95), Inches(6.35) - y, size=12)
            else:
                if s.kind == "video" and embed_video and VIDEO.exists():
                    poster = str(SCREENSHOTS / "region-zh.png")
                    slide.shapes.add_movie(str(VIDEO), Inches(0.6), Inches(1.3), Inches(8.4),
                                           Inches(4.75), poster_frame_image=poster,
                                           mime_type="video/mp4")
                elif img:
                    place(img, Inches(0.6), Inches(1.3), Inches(8.4), Inches(5.0))
                else:
                    text(slide, Inches(0.6), Inches(3.0), Inches(8.4), Inches(1),
                         [[(f"（{s.image} 未產生）", {"color": DIM})]])
                if s.kind == "video" and not embed_video:
                    text(slide, Inches(0.6), Inches(6.05), Inches(8.4), Inches(0.4),
                         [[(f"影片：{VIDEO_LINK}", {"color": ACCENT, "size": 11, "font": MONO})]])
                bullets(slide, s.bullets, Inches(9.2), Inches(1.3), Inches(3.7), Inches(5.2), size=13)
            text(slide, Inches(0.6), Inches(6.45), Inches(12.2), Inches(0.4),
                 [[(s.caption, {"color": DIM, "size": 9.5, "font": MONO})]])
        elif s.kind == "five":
            y = title_line(slide, s, size=24)
            if s.stats:
                y = stat_strip(slide, s.stats, y)
            five_box(slide, s, Inches(0.5), y, Inches(12.3), Inches(6.85) - y)
        chrome(slide, s, number)
        slide.notes_slide.notes_text_frame.text = _notes_text(s)

    target.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(target))
    return target


# ------------------------------------------------------------------ html

def build_html(target: Path, images: dict[str, Path]) -> Path:
    def esc(x: object) -> str:
        return html.escape(str(x), quote=True)

    sections = []
    for number, s in enumerate(SLIDES, 1):
        core = (f'<span class="core">★ 主線 {s.core}/{CORE_TOTAL}</span>' if s.core else "")
        body = ""
        if s.stats:
            body += '<div class="stats">' + "".join(
                f'<div class="stat"><b>{esc(h)}</b><span>{esc(l)}</span></div>' for h, l in s.stats) + "</div>"
        if s.kind == "five":
            body += '<dl class="five">' + "".join(
                f'<dt class="{k}">{esc(v)}</dt><dd>{esc(s.five[k])}</dd>'
                for k, v in FIVE_LABELS.items()) + "</dl>"
        if s.table:
            head, *rows = s.table
            body += ('<div class="scroll"><table><thead><tr>' + "".join(f"<th>{esc(h)}</th>" for h in head)
                     + "</tr></thead><tbody>" + "".join("<tr>" + "".join(f"<td>{esc(c)}</td>" for c in row) + "</tr>" for row in rows)
                     + "</tbody></table></div>")
        if s.bullets:
            body += "<ul>" + "".join(f"<li>{esc(b)}</li>" for b in s.bullets) + "</ul>"
        if s.image and s.image in images:
            body = (f'<figure><img src="img/{esc(s.image)}" alt="{esc(s.caption)}">'
                    f"<figcaption>{esc(s.caption)}</figcaption></figure>") + body
        if s.kind == "video":
            body += (f'<p class="link"><a href="{VIDEO_LINK}">影片（GitHub Release）</a></p>')
        qa = ""
        if s.questions:
            qa = '<div class="qa"><h3>面試官會問</h3>' + "".join(
                f'<div class="q"><p class="question">Q · {esc(q)}</p>'
                f'<p class="answer">{esc(a)}</p></div>' for q, a in s.questions) + "</div>"
        notes = ('<details class="notes"><summary>講稿</summary><p class="plain">'
                 f'{esc(s.plain)}</p><ol>' + "".join(f"<li>{esc(n)}</li>" for n in s.notes)
                 + "</ol></details>")
        sections.append(
            f'<section id="s{number}" data-layer="{esc(s.layer)}">'
            f'<header><span class="layer">{esc(s.layer)}</span>{core}'
            f'<span class="n">{number} / {len(SLIDES)}</span></header>'
            f"<h2>{esc(s.title)}</h2>{body}{qa}{notes}</section>"
        )

    page = f"""<!doctype html>
<html lang="zh-TW">
<head>
<meta charset="utf-8">
<meta name="viewport" content="width=device-width, initial-scale=1">
<title>AOI 複判：一個九天的實驗紀錄</title>
<link rel="stylesheet" href="https://fonts.googleapis.com/css2?family=Noto+Sans+TC:wght@400;500;700&family=Noto+Serif+TC:wght@700&family=IBM+Plex+Mono:wght@400;500&display=swap">
<style>
:root {{ --ground:#{GROUND}; --panel:#{PANEL}; --ink:#{INK}; --dim:#{DIM}; --accent:#{ACCENT}; --warn:#{WARN}; --rule:#33383D; }}
* {{ box-sizing:border-box; }}
body {{ margin:0; background:var(--ground); color:var(--ink); font-family:"Noto Sans TC",system-ui,sans-serif; line-height:1.6; }}
section {{ min-height:100vh; padding:2.2rem 1.4rem 3rem; max-width:1100px; margin:0 auto; border-bottom:1px solid var(--rule); scroll-margin-top:0; }}
header {{ display:flex; gap:1rem; align-items:baseline; font-family:"IBM Plex Mono",monospace; font-size:.78rem; color:var(--dim); }}
header .core {{ color:var(--accent); }} header .n {{ margin-left:auto; }}
h2 {{ font-family:"Noto Serif TC",serif; font-size:1.6rem; line-height:1.3; margin:.6rem 0 1rem; text-wrap:balance; }}
h3 {{ font-size:.95rem; color:var(--dim); margin:1.4rem 0 .4rem; letter-spacing:.04em; }}
ul {{ padding-left:1.2rem; }} li {{ margin:.35rem 0; max-width:80ch; }}
dl.five {{ display:grid; grid-template-columns:7rem 1fr; gap:.5rem 1rem; margin:0 0 1rem; }}
dl.five dt {{ color:var(--accent); font-weight:700; font-size:.9rem; padding-top:.15rem; }}
dl.five dt.mistake {{ color:var(--warn); }}
dl.five dt.mistake, dl.five dt.mistake + dd {{ background:#2B2416; }}
dl.five dt.rule, dl.five dt.rule + dd {{ background:#16292B; }}
dl.five dt, dl.five dd {{ padding:.35rem .6rem; border-radius:4px; }}
dl.five dd {{ margin:0; max-width:85ch; }}
.scroll {{ overflow-x:auto; margin:0 0 1rem; }}
table {{ border-collapse:collapse; width:100%; font-size:.9rem; }}
th {{ text-align:left; color:var(--accent); font-weight:600; font-size:.8rem; padding:.45rem .6rem; border-bottom:1px solid var(--rule); }}
td {{ padding:.45rem .6rem; border-bottom:1px solid var(--rule); vertical-align:top; }}
figure {{ margin:0 0 1rem; background:var(--panel); border:1px solid var(--rule); border-radius:6px; padding:.6rem; }}
figure img {{ max-width:100%; height:auto; display:block; }}
figcaption {{ font-family:"IBM Plex Mono",monospace; font-size:.72rem; color:var(--dim); margin-top:.4rem; }}
.stats {{ display:grid; grid-template-columns:repeat(auto-fit,minmax(180px,1fr)); gap:.6rem; margin:0 0 1rem; }}
.stat {{ background:var(--panel); border:1px solid var(--rule); border-radius:6px; padding:.6rem .8rem; }}
.stat b {{ display:block; font-family:"IBM Plex Mono",monospace; color:var(--accent); font-size:1.3rem; }}
.stat span {{ font-size:.8rem; color:var(--dim); }}
.qa {{ background:var(--panel); border:1px solid var(--rule); border-radius:6px; padding:.2rem 1rem .8rem; margin-top:1.2rem; }}
.q {{ margin:.7rem 0; }} .question {{ margin:0; font-weight:500; }} .answer {{ margin:.2rem 0 0; color:var(--ink); border-left:3px solid var(--accent); padding-left:.7rem; max-width:80ch; }}
details.notes {{ margin-top:1rem; color:var(--dim); font-size:.92rem; }} details.notes summary {{ cursor:pointer; }}
.plain {{ color:var(--ink); font-weight:500; }}
.link a {{ color:var(--accent); }}
body.quiz dl.five, body.quiz ul, body.quiz figure, body.quiz details.notes {{ display:none; }}
body.quiz .answer {{ visibility:hidden; }} body.quiz.reveal .answer {{ visibility:visible; }}
#hud {{ position:fixed; right:.8rem; bottom:.8rem; font-family:"IBM Plex Mono",monospace; font-size:.72rem; color:var(--dim); background:var(--panel); border:1px solid var(--rule); border-radius:4px; padding:.3rem .6rem; }}
@media (prefers-reduced-motion: reduce) {{ html {{ scroll-behavior:auto; }} }}
</style>
</head>
<body>
{"".join(sections)}
<div id="hud">← → 翻頁 · q 自測 · a 看答案 · n 講稿</div>
<script>
(function () {{
  const secs = [...document.querySelectorAll('section')];
  let at = 0;
  function go(i) {{ at = Math.max(0, Math.min(secs.length - 1, i)); secs[at].scrollIntoView({{behavior:'smooth'}}); }}
  document.addEventListener('keydown', function (e) {{
    if (e.key === 'ArrowRight' || e.key === ' ') {{ e.preventDefault(); go(at + 1); }}
    else if (e.key === 'ArrowLeft') {{ e.preventDefault(); go(at - 1); }}
    else if (e.key === 'q') {{ document.body.classList.toggle('quiz'); document.body.classList.remove('reveal'); }}
    else if (e.key === 'a') {{ document.body.classList.toggle('reveal'); }}
    else if (e.key === 'n') {{ document.querySelectorAll('details.notes').forEach(d => d.open = !d.open); }}
  }});
  const io = new IntersectionObserver(es => es.forEach(x => {{ if (x.isIntersecting) at = secs.indexOf(x.target); }}), {{threshold: 0.5}});
  secs.forEach(s => io.observe(s));
}})();
</script>
</body>
</html>
"""
    target.write_text(page)
    return target


# ---------------------------------------------------------- study guide

def build_study_guide(target: Path) -> Path:
    lines = ["# AOI 複判專案歷程：問答手冊", "",
             "從 `scripts/deck_content.py` 產生；和投影片、網頁版同一份內容。",
             "每一頁：一句白話 → 為什麼做／怎麼設計／量到什麼／錯在哪／規則 → 面試官會問 → 講稿。", ""]
    core = [s.title for s in sorted((s for s in SLIDES if s.core), key=lambda s: s.core)]
    lines += ["## 十分鐘主線（9 頁）", ""] + [f"{i}. {t}" for i, t in enumerate(core, 1)] + [""]
    for layer, slides in by_layer().items():
        if not slides:
            continue
        lines += [f"## {layer}", ""]
        for s in slides:
            lines += [f"### {s.title}", "", f"**白話：**{s.plain}", ""]
            if s.five:
                for k, label in FIVE_LABELS.items():
                    lines += [f"- **{label}：**{s.five[k]}"]
                lines += [""]
            if s.table:
                head, *rows = s.table
                lines += ["| " + " | ".join(head) + " |", "|" + "---|" * len(head)]
                lines += ["| " + " | ".join(r) + " |" for r in rows] + [""]
            if s.bullets:
                lines += [f"- {b}" for b in s.bullets] + [""]
            if s.questions:
                lines += ["**面試官會問：**", ""]
                for q, a in s.questions:
                    lines += [f"- Q：{q}", f"  - A：{a}"]
                lines += [""]
            lines += ["**講稿：**", ""] + [f"{i}. {n}" for i, n in enumerate(s.notes, 1)] + [""]
    target.write_text("\n".join(lines))
    return target


# ------------------------------------------------------------------ main

def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--embed-video", action="store_true",
                        help=f"also write {EMBEDDED_OUT.relative_to(ROOT)} with the demo embedded")
    parser.add_argument("--no-images", action="store_true", help="reuse docs/deck/img as it is")
    args = parser.parse_args()

    try:
        import pptx  # noqa: F401
    except ImportError:
        print("python-pptx is not installed. Run:\n"
              "  uv run --with python-pptx --with playwright python scripts/build_deck.py")
        return 2

    OUT.mkdir(parents=True, exist_ok=True)
    if args.no_images:
        images = {p.name: p for p in IMG.glob("*.png")}
    else:
        images = render_images()
    for name, path in sorted(images.items()):
        print(f"image  {path.relative_to(ROOT)}  {path.stat().st_size // 1024} KB")
    missing = sorted({s.image for s in SLIDES if s.image and s.image not in images})
    if missing:
        print(f"slides name images that do not exist: {missing}\n"
              "A slide whose picture is missing renders blank and says nothing.",
              file=sys.stderr)
        return 2

    pptx_path = build_pptx(OUT / "aoi-agent-journey.zh-TW.pptx", images, embed_video=False)
    print(f"pptx   {pptx_path.relative_to(ROOT)}  {pptx_path.stat().st_size // 1024} KB")
    html_path = build_html(OUT / "aoi-agent-journey.zh-TW.html", images)
    print(f"html   {html_path.relative_to(ROOT)}  {html_path.stat().st_size // 1024} KB")
    guide = build_study_guide(OUT / "study-guide.zh-TW.md")
    print(f"guide  {guide.relative_to(ROOT)}")
    if args.embed_video:
        if not VIDEO.exists():
            print(f"{VIDEO.relative_to(ROOT)} not found: no embedded deck written")
        else:
            EMBEDDED_OUT.parent.mkdir(parents=True, exist_ok=True)
            build_pptx(EMBEDDED_OUT, images, embed_video=True)
            print(f"pptx   {EMBEDDED_OUT.relative_to(ROOT)}  (video embedded, not committed)")
    summary = {"slides": len(SLIDES), "core": CORE_TOTAL,
               "layers": [(name, len(slides)) for name, slides in by_layer().items()]}
    print(json.dumps(summary, ensure_ascii=False))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
